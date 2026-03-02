"""
convert_to_pptx.py

Converts the ChatHealthy Pipeline deck into a PowerPoint file.
Navigation buttons are embedded on every slide.

Usage:
    python convert_to_pptx.py <source_jsx_path> <destination_pptx_path>

The source_jsx_path is accepted for API compatibility but the slide
content is derived from the embedded specification below.

Requirements:
    pip install python-pptx
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy


# ── Colour palette ────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0xDD, 0x16, 0x28)
BG_MID       = RGBColor(0x0D, 0x29, 0x52)
ACCENT       = RGBColor(0x4F, 0xC3, 0xF7)
ACCENT2      = RGBColor(0x26, 0xC6, 0xDA)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE   = RGBColor(0x90, 0xCA, 0xF9)
BODY_TEXT    = RGBColor(0xB0, 0xBE, 0xC5)
MUTED        = RGBColor(0x54, 0x6E, 0x7A)
CARD_BORDER  = RGBColor(0x1A, 0x3A, 0x5C)
CARD_BG      = RGBColor(0x0D, 0x1F, 0x3C)

# Layer colours for architecture slide
LAYER_COLORS = {
    "Scheduler":              RGBColor(0x54, 0x6E, 0x7A),
    "Fetch Layer":            RGBColor(0x15, 0x65, 0xC0),
    "Storage":                RGBColor(0x00, 0x69, 0x5C),
    "Processing Layer":       RGBColor(0x6A, 0x1B, 0x9A),
    "Enrichment & Correction":RGBColor(0xC6, 0x28, 0x28),
    "Output":                 RGBColor(0xE6, 0x51, 0x00),
}

STEP_COLORS = {
    "trigger":  RGBColor(0x15, 0x65, 0xC0),
    "check":    RGBColor(0x54, 0x6E, 0x7A),
    "action":   RGBColor(0x00, 0x69, 0x5C),
    "validate": RGBColor(0xE6, 0x51, 0x00),
    "db":       RGBColor(0x6A, 0x1B, 0x9A),
    "enrich":   RGBColor(0xC6, 0x28, 0x28),
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SLIDE_TITLES = [
    "Cover",
    "Problem Statement",
    "Component Architecture",
    "Sequence: Fetch & Index",
    "Sequence: Parallel Processing",
    "Requirements: Fetch & Index",
    "Requirements: Processing",
    "Requirements: Infrastructure",
]


# ── Helpers ───────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_name="Georgia", font_size=12, bold=False,
                 color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_label(slide, text, left, top, width, height,
              font_size=8, color=MUTED, bold=False, align=PP_ALIGN.LEFT):
    return add_text_box(slide, text, left, top, width, height,
                        font_name="Courier New", font_size=font_size,
                        bold=bold, color=color, align=align)


def add_slide_header(slide, eyebrow: str, title: str):
    """Standard header used on all content slides."""
    add_text_box(slide, eyebrow,
                 Inches(0.55), Inches(0.28), Inches(10), Inches(0.25),
                 font_name="Courier New", font_size=8, color=ACCENT,
                 align=PP_ALIGN.LEFT)
    add_text_box(slide, title,
                 Inches(0.55), Inches(0.52), Inches(10), Inches(0.6),
                 font_name="Georgia", font_size=28, bold=True,
                 color=WHITE, align=PP_ALIGN.LEFT)
    # Accent rule
    add_rect(slide, Inches(0.55), Inches(1.18), Inches(0.7), Pt(2),
             fill_color=ACCENT)


def add_card(slide, title, bullets, left, top, width, height):
    """Add a requirement card with title and bullet points."""
    # Card background
    add_rect(slide, left, top, width, height,
             fill_color=CARD_BG, line_color=CARD_BORDER, line_width=Pt(0.75))
    # Card title bar
    add_rect(slide, left, top, width, Inches(0.28),
             fill_color=RGBColor(0x0A, 0x2A, 0x4A))
    add_text_box(slide, title,
                 left + Inches(0.12), top + Inches(0.03),
                 width - Inches(0.2), Inches(0.25),
                 font_name="Courier New", font_size=7.5, bold=True,
                 color=ACCENT, align=PP_ALIGN.LEFT)
    # Bullets
    bullet_top = top + Inches(0.33)
    bullet_h = (height - Inches(0.38)) / max(len(bullets), 1)
    for b in bullets:
        # Triangle marker
        add_text_box(slide, "▸",
                     left + Inches(0.08), bullet_top,
                     Inches(0.15), bullet_h,
                     font_name="Georgia", font_size=7, color=ACCENT)
        add_text_box(slide, b,
                     left + Inches(0.22), bullet_top,
                     width - Inches(0.32), bullet_h,
                     font_name="Georgia", font_size=7, color=BODY_TEXT,
                     wrap=True)
        bullet_top += bullet_h


def add_nav_buttons(slide, slide_index, total, prs):
    """
    Add Prev / slide-number dots / Next navigation to every slide.
    Buttons use hyperlink actions to jump between slides.
    """
    btn_w = Inches(0.7)
    btn_h = Inches(0.28)
    btn_top = SLIDE_H - Inches(0.42)

    # ── PREV button ──────────────────────────────────────────────────────────
    if slide_index > 0:
        prev = slide.shapes.add_shape(1,
            Inches(0.3), btn_top, btn_w, btn_h)
        prev.fill.solid()
        prev.fill.fore_color.rgb = RGBColor(0x0D, 0x29, 0x52)
        prev.line.color.rgb = ACCENT
        prev.line.width = Pt(0.5)
        tf = prev.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = "◀  PREV"
        run.font.name = "Courier New"
        run.font.size = Pt(7)
        run.font.color.rgb = ACCENT
        # Hyperlink to previous slide
        _add_slide_hyperlink(prev, slide_index - 1, prs)

    # ── slide number dots ────────────────────────────────────────────────────
    dot_area_left = Inches(1.15)
    dot_spacing   = (SLIDE_W - Inches(2.3)) / total
    for i in range(total):
        cx = dot_area_left + dot_spacing * i + dot_spacing / 2
        cy = btn_top + btn_h / 2
        r  = Inches(0.055) if i == slide_index else Inches(0.04)
        dot = slide.shapes.add_shape(9,   # oval
            cx - r, cy - r, r * 2, r * 2)
        if i == slide_index:
            dot.fill.solid()
            dot.fill.fore_color.rgb = ACCENT
            dot.line.fill.background()
        else:
            dot.fill.solid()
            dot.fill.fore_color.rgb = MUTED
            dot.line.fill.background()
        # Make each dot a hyperlink to its slide
        _add_slide_hyperlink(dot, i, prs)

    # ── NEXT button ──────────────────────────────────────────────────────────
    if slide_index < total - 1:
        nxt = slide.shapes.add_shape(1,
            SLIDE_W - Inches(1.0), btn_top, btn_w, btn_h)
        nxt.fill.solid()
        nxt.fill.fore_color.rgb = RGBColor(0x0D, 0x29, 0x52)
        nxt.line.color.rgb = ACCENT
        nxt.line.width = Pt(0.5)
        tf = nxt.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = "NEXT  ▶"
        run.font.name = "Courier New"
        run.font.size = Pt(7)
        run.font.color.rgb = ACCENT
        _add_slide_hyperlink(nxt, slide_index + 1, prs)


def _add_slide_hyperlink(shape, target_slide_index, prs):
    """Attach a click-to-slide hyperlink action to a shape.
    target_slide_index is 0-based; python-pptx target_slide expects a Slide object.
    """
    click = shape.click_action
    target_slide = prs.slides[target_slide_index]
    click.target_slide = target_slide


# ── Slide builders ────────────────────────────────────────────────────────────

def build_cover(slide):
    set_bg(slide, BG_DARK)

    # Company name
    add_text_box(slide, "Chat",
                 Inches(3.5), Inches(2.0), Inches(3), Inches(0.85),
                 font_name="Georgia", font_size=52, bold=True,
                 color=WHITE, align=PP_ALIGN.RIGHT)
    add_text_box(slide, "Healthy",
                 Inches(6.4), Inches(2.0), Inches(3.5), Inches(0.85),
                 font_name="Georgia", font_size=52, bold=True,
                 color=ACCENT, align=PP_ALIGN.LEFT)

    # Tagline
    add_text_box(slide, "DATA  INTELLIGENCE  PLATFORM",
                 Inches(2.5), Inches(2.95), Inches(8.3), Inches(0.38),
                 font_name="Courier New", font_size=11, bold=False,
                 color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # Divider rule
    add_rect(slide, Inches(5.4), Inches(3.42), Inches(2.5), Pt(1.5),
             fill_color=ACCENT)

    # Subtitle
    add_text_box(slide, "National Provider Data Pipeline",
                 Inches(2.0), Inches(3.6), Inches(9.3), Inches(0.42),
                 font_name="Georgia", font_size=18, bold=False,
                 color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Architecture & Technical Specification",
                 Inches(2.0), Inches(4.08), Inches(9.3), Inches(0.35),
                 font_name="Courier New", font_size=11, bold=False,
                 color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # SVG-style decorative heart + nodes (drawn with shapes)
    # Outer circle (glow stand-in)
    add_rect(slide, Inches(5.9), Inches(0.6), Inches(1.5), Inches(1.5),
             fill_color=RGBColor(0x0A, 0x22, 0x40))
    # Three pipeline nodes
    for cx, color in [(Inches(6.1), ACCENT), (Inches(6.65), ACCENT2), (Inches(7.2), ACCENT)]:
        r = Inches(0.1)
        dot = slide.shapes.add_shape(9, cx, Inches(1.15), r*2, r*2)
        dot.fill.solid(); dot.fill.fore_color.rgb = color
        dot.line.fill.background()

    # Date / footer
    add_text_box(slide, "March 2026  |  CONFIDENTIAL",
                 Inches(0.4), Inches(7.1), Inches(5), Inches(0.25),
                 font_name="Courier New", font_size=8, color=MUTED)
    add_text_box(slide, "chathealthy.com",
                 Inches(8.3), Inches(7.1), Inches(4.6), Inches(0.25),
                 font_name="Courier New", font_size=8, color=MUTED,
                 align=PP_ALIGN.RIGHT)


def build_problem(slide):
    set_bg(slide, BG_DARK)
    add_slide_header(slide, "THE CHALLENGE", "Problem Statement")

    problems = [
        ("⚕  Stale Provider Data",
         "The NPPES National Provider file — 15M+ records — is published monthly by CMS. "
         "Without automation, data quickly becomes outdated, leading to incorrect provider "
         "directories and compliance risk."),
        ("⚙  Manual Processing Bottleneck",
         "The full file is 8–12 GB uncompressed. Manual download, validation, and ingestion "
         "is error-prone, slow, and unsustainable as data volume grows."),
        ("🔎  Raw Data is Incomplete",
         "CMS source data contains known errors: malformed addresses, missing taxonomy codes, "
         "inconsistent phone formats, and encoding issues. Raw data cannot be trusted without "
         "enrichment and correction."),
        ("📊  No Scalable Processing Model",
         "Sequential processing of 15M records is impractical. A scalable parallel architecture "
         "is required to meet turnaround SLAs and support downstream consumers."),
    ]

    col_w = Inches(5.9)
    col_h = Inches(2.3)
    positions = [
        (Inches(0.45), Inches(1.35)),
        (Inches(6.55), Inches(1.35)),
        (Inches(0.45), Inches(3.75)),
        (Inches(6.55), Inches(3.75)),
    ]

    for (left, top), (title, body) in zip(positions, problems):
        add_rect(slide, left, top, col_w, col_h,
                 fill_color=CARD_BG, line_color=CARD_BORDER, line_width=Pt(0.75))
        add_text_box(slide, title,
                     left + Inches(0.15), top + Inches(0.12),
                     col_w - Inches(0.25), Inches(0.35),
                     font_name="Georgia", font_size=12, bold=True, color=ACCENT)
        add_text_box(slide, body,
                     left + Inches(0.15), top + Inches(0.52),
                     col_w - Inches(0.25), Inches(1.65),
                     font_name="Georgia", font_size=9.5, color=BODY_TEXT, wrap=True)

    # Goal bar
    add_rect(slide, Inches(0.45), Inches(6.2), Inches(12.43), Inches(0.58),
             fill_color=RGBColor(0x08, 0x1E, 0x3A), line_color=ACCENT, line_width=Pt(0.75))
    add_text_box(slide,
                 "Goal:  Build a fully automated, cloud-native pipeline on Azure that detects, "
                 "downloads, validates, enriches, and processes the NPPES monthly file — reliably and at scale.",
                 Inches(0.6), Inches(6.26), Inches(12.1), Inches(0.46),
                 font_name="Georgia", font_size=9.5, color=WHITE,
                 align=PP_ALIGN.CENTER, wrap=True)


def build_architecture(slide):
    set_bg(slide, BG_DARK)
    add_slide_header(slide, "SYSTEM DESIGN", "Component Architecture")

    layers = [
        ("Scheduler",               ["Azure Logic App", "Monthly trigger", "Offset from month-end"]),
        ("Fetch Layer",             ["FileWatcher", "FileDownloader (streaming)", "FileIndexer", "Azure Container Instance"]),
        ("Storage",                 ["Azure Blob — Cool tier", "3-month retention", "90-day lifecycle policy", "MongoDB — chunk index"]),
        ("Processing Layer",        ["ChunkProcessor (fan-out)", "150 parallel workers", "Azure Container Apps", "Strategy Pattern invokers"]),
        ("Enrichment & Correction", ["NUCC Taxonomy lookup", "Address standardization", "Phone normalization", "Error rules engine"]),
        ("Output",                  ["MongoDB Atlas (enriched)", "Audit log per chunk", "Downstream API consumers"]),
    ]

    n = len(layers)
    col_w = (SLIDE_W - Inches(1.0)) / n
    col_top = Inches(1.35)
    col_h = Inches(4.9)

    for i, (layer_name, items) in enumerate(layers):
        left = Inches(0.5) + col_w * i
        color = LAYER_COLORS[layer_name]

        # Header bar
        add_rect(slide, left, col_top, col_w - Inches(0.08), Inches(0.3),
                 fill_color=color)
        add_text_box(slide, layer_name.upper(),
                     left + Inches(0.04), col_top + Inches(0.04),
                     col_w - Inches(0.14), Inches(0.22),
                     font_name="Courier New", font_size=6.5, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER)

        # Body
        body_top = col_top + Inches(0.3)
        body_h = col_h - Inches(0.3)
        add_rect(slide, left, body_top, col_w - Inches(0.08), body_h,
                 fill_color=CARD_BG, line_color=color, line_width=Pt(0.5))

        item_top = body_top + Inches(0.12)
        for item in items:
            add_rect(slide, left + Inches(0.06), item_top,
                     col_w - Inches(0.2), Inches(0.5),
                     fill_color=RGBColor(0x0A, 0x16, 0x28))
            add_text_box(slide, item,
                         left + Inches(0.1), item_top + Inches(0.06),
                         col_w - Inches(0.26), Inches(0.42),
                         font_name="Courier New", font_size=7.5, color=BODY_TEXT, wrap=True)
            item_top += Inches(0.6)

    # Flow line
    flow = "Trigger  →  Fetch  →  Store + Index  →  Fan-out  →  Enrich  →  Output"
    add_text_box(slide, flow,
                 Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.25),
                 font_name="Courier New", font_size=8, color=ACCENT,
                 align=PP_ALIGN.CENTER)


def build_sequence(slide, eyebrow, title, steps):
    """
    steps = list of (actor, action, type_key)
    """
    set_bg(slide, BG_DARK)
    add_slide_header(slide, eyebrow, title)

    row_h = Inches(0.4)
    top = Inches(1.38)
    num_w = Inches(0.28)
    actor_w = Inches(1.8)
    action_left = Inches(0.5) + num_w + actor_w + Inches(0.1)
    action_w = SLIDE_W - action_left - Inches(0.4)

    for i, (actor, action, type_key) in enumerate(steps):
        color = STEP_COLORS.get(type_key, MUTED)
        row_top = top + row_h * i

        # Number badge
        badge = slide.shapes.add_shape(9,
            Inches(0.5), row_top + Inches(0.05), Inches(0.24), Inches(0.24))
        badge.fill.solid(); badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        tf = badge.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = str(i + 1)
        run.font.name = "Courier New"; run.font.size = Pt(7)
        run.font.color.rgb = WHITE; run.font.bold = True

        # Actor label
        add_text_box(slide, actor,
                     Inches(0.5) + num_w + Inches(0.06), row_top,
                     actor_w, row_h,
                     font_name="Courier New", font_size=8, color=ACCENT)

        # Left border accent
        add_rect(slide, action_left - Inches(0.06), row_top + Inches(0.06),
                 Pt(2), row_h - Inches(0.12), fill_color=color)

        # Action text
        add_text_box(slide, action,
                     action_left, row_top, action_w, row_h,
                     font_name="Georgia", font_size=8.5, color=BODY_TEXT, wrap=True)

    # Legend
    legend_items = list(STEP_COLORS.items())
    legend_labels = {
        "trigger": "Trigger", "check": "Decision",
        "action": "Process",  "validate": "Validate",
        "db": "MongoDB",       "enrich": "Enrich"
    }
    legend_top = SLIDE_H - Inches(0.55)
    lx = Inches(0.5)
    for key, color in legend_items:
        dot = slide.shapes.add_shape(9, lx, legend_top + Inches(0.06),
                                     Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        add_text_box(slide, legend_labels[key],
                     lx + Inches(0.14), legend_top,
                     Inches(0.8), Inches(0.22),
                     font_name="Courier New", font_size=7, color=MUTED)
        lx += Inches(1.05)


def build_requirements(slide, eyebrow, title, cards):
    """
    cards = list of (card_title, [bullet, ...])
    Laid out in a 2x2 grid.
    """
    set_bg(slide, BG_DARK)
    add_slide_header(slide, eyebrow, title)

    n_cols = 2
    margin_l = Inches(0.45)
    margin_t = Inches(1.38)
    gap = Inches(0.12)
    usable_w = SLIDE_W - margin_l * 2
    usable_h = SLIDE_H - margin_t - Inches(0.65)
    card_w = (usable_w - gap) / n_cols
    card_h = (usable_h - gap) / 2

    for idx, (card_title, bullets) in enumerate(cards):
        col = idx % n_cols
        row = idx // n_cols
        left = margin_l + col * (card_w + gap)
        top  = margin_t + row * (card_h + gap)
        add_card(slide, card_title, bullets, left, top, card_w, card_h)


# ── Main builder ──────────────────────────────────────────────────────────────

def build_presentation(source_path: str, dest_path: str):
    """
    Build the ChatHealthy pipeline deck and save to dest_path.
    source_path is accepted for API compatibility but not parsed.
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # completely blank layout

    # We need all slides added first so hyperlinks can reference them.
    all_slides = [prs.slides.add_slide(blank_layout) for _ in SLIDE_TITLES]

    # ── Slide 0: Cover ────────────────────────────────────────────────────────
    build_cover(all_slides[0])

    # ── Slide 1: Problem ──────────────────────────────────────────────────────
    build_problem(all_slides[1])

    # ── Slide 2: Architecture ─────────────────────────────────────────────────
    build_architecture(all_slides[2])

    # ── Slide 3: Sequence — Fetch & Index ─────────────────────────────────────
    fetch_steps = [
        ("Logic App",        "Trigger monthly job (offset 3 days from month-end)", "trigger"),
        ("FileWatcher",      "HTTP GET cms.gov download page — parse filename for month/year", "check"),
        ("FileWatcher",      "Compare filename + SHA256 against file_registry in MongoDB — already processed?", "check"),
        ("FileDownloader",   "Stream HTTPS download from CMS → unzip transform → Azure Blob multipart upload (never writes to disk)", "action"),
        ("FileDownloader",   "Validate SHA256 checksum of uploaded CSV", "validate"),
        ("file_registry",    "Insert document: file_name, file_url, file_month, total_bytes, checksum, status='arrived'", "db"),
        ("FileIndexer",      "Single sequential read of CSV stream — record byte offset of every 100,000th line", "action"),
        ("file_chunks",      "Bulk insert 150 chunk documents: byte_start, byte_end, record_start, record_end, status='pending'", "db"),
        ("file_registry",    "Update status → 'indexed'", "db"),
        ("Logic App",        "Trigger ChunkProcessor fan-out", "trigger"),
    ]
    build_sequence(all_slides[3], "SEQUENCE 1 OF 2", "Fetch & Index Flow", fetch_steps)

    # ── Slide 4: Sequence — Parallel Processing ───────────────────────────────
    process_steps = [
        ("ChunkProcessor",   "Query MongoDB: find all file_chunks WHERE file_id=X AND status='pending'", "db"),
        ("ChunkProcessor",   "Fan-out: spawn 150 parallel Azure Container App worker instances", "trigger"),
        ("Worker (each)",    "Atomic findOneAndUpdate in MongoDB — claim chunk, set status='processing', record worker_id", "db"),
        ("Worker (each)",    "HTTP Range GET against Azure Blob — retrieve only assigned byte range (no full file read)", "action"),
        ("Worker (each)",    "Trim partial first/last lines at byte boundaries", "action"),
        ("NPPESInvoker",     "Invoke NPPES microservice via Strategy Pattern — enrichment + error correction on 100K records", "enrich"),
        ("NPPESInvoker",     "Enrichment: NUCC taxonomy labels, address standardization, phone normalization", "enrich"),
        ("NPPESInvoker",     "Error correction: encoding fixes, missing field defaults, deduplication check", "validate"),
        ("Worker (each)",    "Write enriched records to MongoDB output collection", "db"),
        ("file_chunks",      "Update chunk status → 'complete' or 'failed' with error detail", "db"),
        ("ChunkProcessor",   "Monitor progress — auto-retry failed chunks up to 3 attempts with exponential backoff", "check"),
        ("file_registry",    "Update status → 'complete' when all chunks finished. Emit completion event.", "db"),
    ]
    build_sequence(all_slides[4], "SEQUENCE 2 OF 2", "Parallel Processing Flow", process_steps)

    # ── Slide 5: Requirements — Fetch & Index ─────────────────────────────────
    build_requirements(all_slides[5],
        "REQUIREMENTS — PAGE 1 OF 3", "Fetch & Index Service",
        [
            ("FileWatcher Class", [
                "HTTP GET the CMS NPPES download page at https://download.cms.gov/nppes/NPI_Files.html",
                "Parse HTML to extract current monthly full-file filename and derive release month/year",
                "Query MongoDB file_registry for a document matching file_month — if found and status is not 'failed', abort with no-op log",
                "Compute SHA256 of remote file via HTTP HEAD or partial GET — compare to stored checksum before downloading",
                "Expose a single public method: detect() returning { isNew: boolean, fileUrl: string, fileName: string }",
            ]),
            ("FileDownloader Class", [
                "Accept fileUrl and target Azure Blob container/path as constructor arguments",
                "Stream download via HTTPS without writing to local disk at any point",
                "Pipe download stream through a zlib/unzip transform stream to decompress on the fly",
                "Pipe decompressed stream to Azure Blob Storage multipart upload using BlockBlobClient.uploadStream()",
                "On completion, compute and return SHA256 checksum of the uploaded CSV",
                "Handle failure: if upload stream fails, delete partial blob and restart from byte 0 — no partial resume",
                "Emit progress events: bytes downloaded, bytes uploaded, percentage complete",
            ]),
            ("FileIndexer Class", [
                "Accept Azure Blob CSV URL and MongoDB connection as constructor arguments",
                "Stream-read the CSV from Azure Blob in sequential chunks — never load full file into memory",
                "Count line returns to build byte-offset index: record byte position of line 1, 100001, 200001, etc.",
                "Chunk size is configurable — default 100,000 records per chunk",
                "Bulk insert all chunk documents into MongoDB file_chunks collection in a single operation",
                "Each document must contain: file_id, chunk_number, byte_start, byte_end, record_start, record_end, status='pending'",
                "Update file_registry status to 'indexed' and set total_records and total_chunks on completion",
            ]),
            ("MongoDB Schema — file_registry", [
                "Fields: _id (ObjectId), file_name (string), file_url (string), file_month (Date), total_records (long), total_bytes (long), total_chunks (int), checksum (string), indexed_at (Date), status (enum: arrived|indexed|processing|complete|failed)",
                "Unique index on file_month to prevent duplicate processing",
                "Index on status for pipeline monitoring queries",
            ]),
        ]
    )

    # ── Slide 6: Requirements — Processing ───────────────────────────────────
    build_requirements(all_slides[6],
        "REQUIREMENTS — PAGE 2 OF 3", "Processing & Enrichment Service",
        [
            ("ChunkProcessor Class", [
                "Query MongoDB file_chunks for all pending chunks for a given file_id",
                "Fan-out: invoke one Azure Container App worker instance per chunk — pass chunk _id only",
                "Each worker uses MongoDB findOneAndUpdate with status='pending' filter to atomically claim its chunk",
                "Monitor chunk statuses via polling or MongoDB change streams",
                "Auto-retry any chunk with status='failed' up to 3 times with exponential backoff",
                "On all chunks complete, update file_registry status to 'complete'",
                "Accept a processorStrategy (InvokerInterface) via constructor — never reference a specific file type internally",
            ]),
            ("InvokerInterface Contract", [
                "All microservice invokers must implement: invoke(chunkMeta), onSuccess(chunkMeta, result), onFailure(chunkMeta, error)",
                "chunkMeta object contains: file_id, chunk_number, byte_start, byte_end, record_start, record_end, file_url",
                "invoke() must perform HTTP Range GET against Azure Blob using byte_start and byte_end — header: Range: bytes={start}-{end}",
                "invoke() must handle partial first and last lines: discard bytes before first newline at start, read through first newline past end byte",
                "invoke() must return enriched records array and a corrections_log array",
            ]),
            ("NPPESInvoker — Enrichment Rules", [
                "Taxonomy code lookup: join each taxonomy_code against NUCC crosswalk to append taxonomy_description and taxonomy_grouping",
                "Address standardization: normalize state abbreviations, fix ZIP code formatting, remove non-printable characters",
                "Phone normalization: strip non-numeric characters, validate 10-digit US format, null out invalid values",
                "Encoding: convert all fields to UTF-8, replace or remove invalid byte sequences",
                "Flag records missing both primary address and mailing address as status='incomplete'",
            ]),
            ("MongoDB Schema — file_chunks", [
                "Fields: _id (ObjectId), file_id (ObjectId ref), chunk_number (int), byte_start (long), byte_end (long), record_start (long), record_end (long), status (enum: pending|processing|complete|failed), worker_id (string), started_at (Date), completed_at (Date), error (string), retry_count (int)",
                "Compound index on { file_id: 1, status: 1 } for fan-out queries",
                "Compound index on { file_id: 1, chunk_number: 1 } for direct chunk lookup",
            ]),
        ]
    )

    # ── Slide 7: Requirements — Infrastructure ────────────────────────────────
    build_requirements(all_slides[7],
        "REQUIREMENTS — PAGE 3 OF 3", "Infrastructure & Operational Requirements",
        [
            ("Azure Infrastructure", [
                "Storage: Azure Blob Storage Cool tier — container 'nppes-source', 90-day lifecycle delete policy enabled",
                "Fetch compute: Azure Container Instance — no execution time limit, minimum 4 vCPU / 8 GB RAM for streaming",
                "Worker compute: Azure Container Apps — serverless scale-to-zero, max 200 instances, 1 vCPU / 2 GB per worker",
                "Orchestration: Azure Logic Apps — monthly schedule trigger, offset configurable via environment variable",
                "Secrets: Azure Key Vault for MongoDB connection string, Azure Storage connection string, CMS URL",
            ]),
            ("Non-Functional Requirements", [
                "Fetch service must complete download and indexing within 4 hours of trigger",
                "All 150 chunks must complete processing within 6 hours of fan-out trigger",
                "No data may be written to local disk at any stage of the pipeline",
                "All pipeline stages must be idempotent — safe to re-run without duplicating data",
                "Pipeline must detect and skip months where the source file has not changed (checksum match)",
                "All corrections applied to records must be logged with: field_name, original_value, corrected_value, rule_applied",
            ]),
            ("Error Handling Requirements", [
                "FileWatcher: if CMS page is unreachable, retry 3 times at 15-minute intervals before alerting",
                "FileDownloader: if upload stream fails, delete partial blob and restart from byte 0",
                "FileIndexer: if indexing fails mid-way, delete all file_chunks for this file_id and restart",
                "ChunkProcessor: failed chunks retried up to 3 times — after 3 failures mark as 'dead' and raise alert",
                "All errors must be written to MongoDB with timestamp, component name, error message, and stack trace",
            ]),
            ("Schema Change Detection", [
                "At ingest time, validate CSV header row against expected schema document stored in MongoDB pipeline_config collection",
                "If column count, names, or order differ from expected, halt pipeline immediately and raise a critical alert — do not process",
                "Store the detected schema with each file_registry document for audit purposes",
                "Schema expectations must be configurable without code changes — stored in MongoDB pipeline_config collection",
            ]),
        ]
    )

    # ── Navigation buttons on every slide ─────────────────────────────────────
    total = len(all_slides)
    for i, slide in enumerate(all_slides):
        add_nav_buttons(slide, i, total, prs)

    prs.save(dest_path)
    print(f"Saved: {dest_path}")


# ── CLI entry point ────────────────────────────────────────────────────────────

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python convert_to_pptx.py <source_jsx_path> <destination_pptx_path>")
        sys.exit(1)

    source = sys.argv[1]
    dest   = sys.argv[2]

    build_presentation(source, dest)